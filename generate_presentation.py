import os
try:
    from pptx import Presentation
except ImportError:
    print("python-pptx is not installed. Please run 'pip install python-pptx' first.")
    exit(1)

prs = Presentation()

# Standard Layout Definitions
title_slide_layout = prs.slide_layouts[0]
bullet_slide_layout = prs.slide_layouts[1]
two_content_layout = prs.slide_layouts[3]

# --- Slide 1: Introduction ---
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Sudoku Ascent"
subtitle.text = "Project Architecture & Implementation\nTransforming Traditional Sudoku with Modern UI and Dynamic AI"

# --- Slide 2: Methodology & Architecture ---
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
body = slide.placeholders[1].text_frame
title.text = "Methodology & Architecture"
body.text = "We adopted a modular, decoupled microservices approach:"
p = body.add_paragraph()
p.text = "Web Frontend: React & Vite, emphasizing premium glassmorphism and real-time cascaded animations"
p.level = 1
p = body.add_paragraph()
p.text = "Core Backend: Java Spring Boot powers strict puzzle generation rulesets and state logic"
p.level = 1
p = body.add_paragraph()
p.text = "Intelligence Service: Python FastAPI dynamically predicts scaling and ranks generated board mechanics"
p.level = 1

# --- Slide 3: Results & Outcomes ---
slide = prs.slides.add_slide(two_content_layout)
title = slide.shapes.title
title.text = "Results & Outcomes"
left_body = slide.placeholders[1].text_frame
right_body = slide.placeholders[2].text_frame

left_body.text = "Traditional Approach"
p = left_body.add_paragraph()
p.text = "Static, hardcoded difficulty tiers (Easy, Medium, Hard)"
p.level = 1
p = left_body.add_paragraph()
p.text = "Basic, utilitarian user interfaces"
p.level = 1

right_body.text = "Sudoku Ascent"
p = right_body.add_paragraph()
p.text = "AI 'Survival Mode' smoothly scales difficulty strictly by decimal fractions"
p.level = 1
p = right_body.add_paragraph()
p.text = "Premium, reactive UX (e.g., locking 2.8s cascading green win animation validation)"
p.level = 1

# --- Slide 4: Real-time Validation Demo ---
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
body = slide.placeholders[1].text_frame
title.text = "Interactive Application Demo"
body.text = "Showcasing live UI implementations:"
p = body.add_paragraph()
p.text = "Survival Mode intelligent mathematical fetch"
p.level = 1
p = body.add_paragraph()
p.text = "Real-time 'Mistake Checking' matrix validation tool"
p.level = 1
p = body.add_paragraph()
p.text = "Asynchronous algorithmic processing screen"
p.level = 1

# --- Slide 5: Discussion & Challenges ---
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
body = slide.placeholders[1].text_frame
title.text = "Discussion & Technical Challenges"
body.text = "Primary Challenge: Generation Bottlenecks"
p = body.add_paragraph()
p.text = "Problem: Creating a true Sudoku board containing exactly ONE unique solver path natively freezes servers for seconds."
p.level = 1
p = body.add_paragraph()
p.text = "Naive backtracking brute-force arrays cause massive thread congestion."
p.level = 1
p = body.add_paragraph()
p.text = "Solution: Integrating hyper-efficient mathematical exact-cover concepts into the Java core pipeline."
p.level = 1

# --- Slide 6: Code Review (Algorithm X) ---
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
body = slide.placeholders[1].text_frame
title.text = "Code Review: Algorithmic Deep Dive"
body.text = "Donald Knuth's Algorithm X via Dancing Links (DLX)"
p = body.add_paragraph()
p.text = "Re-frames native Sudoku as an 'Exact Cover' problem"
p.level = 1
p = body.add_paragraph()
p.text = "Transforms 9x9 constraints strictly into a sparse 324-column binary matrix"
p.level = 1
p = body.add_paragraph()
p.text = "The DLX mechanism navigates this matrix via node un-hooking instead of array duplication, dropping solve-time to sub-milliseconds."
p.level = 1

# --- Slide 7: Code Review (DLX Node Structure) ---
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
body = slide.placeholders[1].text_frame
title.text = "Code Snippet: The Node Matrix"
body.text = "SudokuSolverService.java - Node Initializer"
p = body.add_paragraph()
p.text = "class Node {\n    Node left, right, up, down;\n    ColumnNode column;\n\n    public Node() {\n        left = right = up = down = this;\n    }\n}"
p = body.add_paragraph()
p.text = "Every '1' in the exact cover matrix is mapped tightly to a custom spatial Node linked specifically to its 4 cardinal neighbors."
p.level = 1

# --- Slide 8: Code Review (Cover Operations) ---
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
body = slide.placeholders[1].text_frame
title.text = "Code Snippet: O(1) Matrix Traversal"
body.text = "SudokuSolverService.java - The DLX 'Cover' Method"
p = body.add_paragraph()
p.text = "void cover() {\n    right.left = left;\n    left.right = right;\n    // ... cascade vertically\n}"
p = body.add_paragraph()
p.text = "Nodes are literally pointer un-linked sequentially to physically 'hide' constraints temporarily during the backtrack."
p.level = 1
p = body.add_paragraph()
p.text = "Re-applying constraints for an incorrect pathway simply reverses the variable reassignment (no array rebuilding required)."
p.level = 1

# --- Slide 9: Conclusion ---
slide = prs.slides.add_slide(bullet_slide_layout)
title = slide.shapes.title
body = slide.placeholders[1].text_frame
title.text = "Conclusion & Future Horizon"
body.text = "Development Summary"
p = body.add_paragraph()
p.text = "Success in obfuscating heavy Java algorithms directly behind an instantaneous React UX"
p.level = 1
p = body.add_paragraph()
p.text = "AI prediction curves implemented without sacrificing organic UI feedback loops"
p.level = 1
p = body.add_paragraph()
p.text = "Future Project Extrapolations:"
p.level = 0
p = body.add_paragraph()
p.text = "Implementing 'Intelligent Hinting' systems by deriving logic directly from DLX pathways"
p.level = 1
p = body.add_paragraph()
p.text = "Transitioning to a cross-platform React Native mobile release"
p.level = 1

output_path = os.path.join(os.getcwd(), "SudokuAscent_Presentation.pptx")
prs.save(output_path)
print(f"Presentation generated successfully at: {output_path}")
